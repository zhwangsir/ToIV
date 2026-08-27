"""文档上传与长文本理解测试(services/docs + routes/documents + runner 注入)。

覆盖:
  - kind_from_filename 扩展名判定
  - split_chunks 段落聚合 / 超长段落硬切 / MAX_CHUNKS 截断
  - parse_text txt/md/docx(pdf 解析依赖真实排版文本,不在单测覆盖)
  - upload/list/delete 端点全流程(embedding mock)
  - embedding 不可用 → 文档保留 + status=no_embed
  - 用户隔离:他人文档不可见/不可删/不可挂载检索
  - runner._docs_context 检索注入(embedding mock)
"""
from __future__ import annotations

import io
import json

import pytest
from fastapi.testclient import TestClient
from sqlalchemy.pool import StaticPool
from sqlmodel import Session, SQLModel, create_engine

from app.agent import rag, runner
from app.db import get_session
from app.main import app
from app.models import Document, Tenant, User
from app.security import create_token, hash_password
from app.services import docs as docsvc


# --------------------------------------------------------------------------- #
# fixtures
# --------------------------------------------------------------------------- #
def _make_user(session: Session, email: str) -> User:
    tenant = Tenant(name=email.split("@")[0])
    session.add(tenant)
    session.commit()
    session.refresh(tenant)
    user = User(
        email=email,
        hashed_password=hash_password("password1"),
        tenant_id=tenant.id,
    )
    session.add(user)
    session.commit()
    session.refresh(user)
    return user


@pytest.fixture
def ctx(tmp_path, monkeypatch):
    """内存库 + 文档目录隔离到 tmp_path + embedding mock(恒定单位向量)。"""
    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    SQLModel.metadata.create_all(engine)

    def override():
        with Session(engine) as session:
            yield session

    app.dependency_overrides[get_session] = override
    monkeypatch.setattr(docsvc, "docs_root", lambda: tmp_path)

    async def fake_embed(texts: list[str]):
        return [[1.0, 0.0, 0.0] for _ in texts]

    monkeypatch.setattr(rag, "_embed", fake_embed)

    with Session(engine) as s:
        alice = _make_user(s, "alice@toiv.ai")
        bob = _make_user(s, "bob@toiv.ai")
        ids = (alice.id, bob.id)
    yield TestClient(app), create_token(ids[0]), create_token(ids[1]), engine, ids
    app.dependency_overrides.clear()


def _auth(token: str) -> dict:
    return {"Authorization": f"Bearer {token}"}


# --------------------------------------------------------------------------- #
# 纯函数:kind / chunk / parse
# --------------------------------------------------------------------------- #
def test_kind_from_filename():
    assert docsvc.kind_from_filename("报告.PDF") == "pdf"
    assert docsvc.kind_from_filename("a.docx") == "docx"
    assert docsvc.kind_from_filename("a.txt") == "txt"
    assert docsvc.kind_from_filename("a.md") == "md"
    assert docsvc.kind_from_filename("a.exe") is None
    assert docsvc.kind_from_filename("noext") is None


def test_split_chunks_short_text_single_chunk():
    chunks, truncated = docsvc.split_chunks("第一段。\n\n第二段。")
    assert chunks == ["第一段。\n第二段。"]
    assert truncated is False


def test_split_chunks_respects_size_and_paragraphs():
    # 每段 400 字符,CHUNK_CHARS=900 → 每块装 2 段
    text = "\n\n".join("甲" * 400 for _ in range(6))
    chunks, truncated = docsvc.split_chunks(text)
    assert len(chunks) == 3
    assert all(len(c) <= docsvc.CHUNK_CHARS for c in chunks)
    assert truncated is False


def test_split_chunks_hard_splits_long_paragraph_with_overlap():
    text = "长" * (docsvc.CHUNK_CHARS * 2 + 100)  # 无换行单段
    chunks, _ = docsvc.split_chunks(text)
    assert len(chunks) == 3
    assert len(chunks[0]) == docsvc.CHUNK_CHARS
    # 重叠:第二块开头 = 第一块结尾的尾巴
    assert chunks[1][:10] == chunks[0][-docsvc.CHUNK_OVERLAP :][:10]


def test_split_chunks_truncates_at_max():
    # 每段略小于 CHUNK_CHARS → 每段独占一块;段数 > MAX_CHUNKS 触发截断
    text = "\n\n".join("字" * (docsvc.CHUNK_CHARS - 10) for _ in range(docsvc.MAX_CHUNKS + 5))
    chunks, truncated = docsvc.split_chunks(text)
    assert truncated is True
    assert len(chunks) == docsvc.MAX_CHUNKS


def test_parse_text_txt_and_md():
    assert docsvc.parse_text("txt", "你好\n世界".encode()) == "你好\n世界"
    assert docsvc.parse_text("md", "# 标题".encode()) == "# 标题"
    # 坏字节不炸,按 replace 容错
    assert "" in docsvc.parse_text("txt", b"\xff\xfe bad")


def test_parse_text_docx_roundtrip():
    import docx

    buf = io.BytesIO()
    d = docx.Document()
    d.add_paragraph("第一章 开始")
    d.add_paragraph("正文内容")
    d.save(buf)
    text = docsvc.parse_text("docx", buf.getvalue())
    assert "第一章 开始" in text and "正文内容" in text


def test_parse_text_unsupported_kind():
    with pytest.raises(ValueError):
        docsvc.parse_text("exe", b"x")


# --------------------------------------------------------------------------- #
# 端点:upload / list / delete
# --------------------------------------------------------------------------- #
def test_upload_list_delete_flow(ctx):
    c, alice_token, _, _, _ = ctx
    r = c.post(
        "/api/docs/upload",
        files={"file": ("笔记.txt", "第一段内容。\n\n第二段内容。".encode(), "text/plain")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 201, r.text
    doc = r.json()
    assert doc["filename"] == "笔记.txt"
    assert doc["kind"] == "txt"
    assert doc["chunk_count"] == 1
    assert doc["status"] == "ready"

    r = c.get("/api/docs", headers=_auth(alice_token))
    assert r.status_code == 200
    assert [d["id"] for d in r.json()] == [doc["id"]]

    r = c.delete(f"/api/docs/{doc['id']}", headers=_auth(alice_token))
    assert r.status_code == 200
    assert c.get("/api/docs", headers=_auth(alice_token)).json() == []


def test_upload_rejects_bad_extension(ctx):
    c, alice_token, _, _, _ = ctx
    r = c.post(
        "/api/docs/upload",
        files={"file": ("a.exe", b"MZ", "application/octet-stream")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 400


def test_upload_rejects_oversize(ctx, monkeypatch):
    c, alice_token, _, _, _ = ctx
    monkeypatch.setattr(docsvc, "MAX_FILE_BYTES", 10)
    r = c.post(
        "/api/docs/upload",
        files={"file": ("big.txt", b"x" * 20, "text/plain")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 413


def test_upload_rejects_empty_text(ctx):
    c, alice_token, _, _, _ = ctx
    r = c.post(
        "/api/docs/upload",
        files={"file": ("empty.txt", b"   \n\n  ", "text/plain")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 422


def test_upload_docx_via_endpoint(ctx):
    import docx

    c, alice_token, _, _, _ = ctx
    buf = io.BytesIO()
    d = docx.Document()
    d.add_paragraph("合同条款一:甲方负责付款。")
    d.save(buf)
    r = c.post(
        "/api/docs/upload",
        files={"file": ("合同.docx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 201, r.text
    assert r.json()["status"] == "ready"


def test_upload_keeps_doc_when_embedding_down(ctx, monkeypatch):
    c, alice_token, _, _, _ = ctx

    async def down_embed(texts):
        return None

    monkeypatch.setattr(rag, "_embed", down_embed)
    r = c.post(
        "/api/docs/upload",
        files={"file": ("笔记.txt", "内容".encode(), "text/plain")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 201
    assert r.json()["status"] == "no_embed"
    # 文档仍在列表中(检索降级为空,不丢文档)
    assert len(c.get("/api/docs", headers=_auth(alice_token)).json()) == 1


def test_user_isolation(ctx):
    c, alice_token, bob_token, _, _ = ctx
    r = c.post(
        "/api/docs/upload",
        files={"file": ("私密.txt", "alice 的秘密".encode(), "text/plain")},
        headers=_auth(alice_token),
    )
    doc_id = r.json()["id"]
    # bob 列表为空、删 alice 文档 404
    assert c.get("/api/docs", headers=_auth(bob_token)).json() == []
    assert c.delete(f"/api/docs/{doc_id}", headers=_auth(bob_token)).status_code == 404
    # alice 的文档还在
    assert len(c.get("/api/docs", headers=_auth(alice_token)).json()) == 1


def test_upload_requires_auth(ctx):
    c, _, _, _, _ = ctx
    r = c.post("/api/docs/upload", files={"file": ("a.txt", b"x", "text/plain")})
    assert r.status_code == 401


# --------------------------------------------------------------------------- #
# runner 检索注入
# --------------------------------------------------------------------------- #
async def test_docs_context_injects_relevant_chunks(ctx):
    _, _, _, engine, (alice_id, _) = ctx
    with Session(engine) as s:
        doc = Document(
            tenant_id="t",
            user_id=alice_id,
            filename="世界观设定.md",
            kind="md",
            size=100,
            chunk_count=1,
        )
        s.add(doc)
        s.commit()
        s.refresh(doc)
        ok = await docsvc.build_index(alice_id, doc.id, ["主角名叫阿澈,是流浪剑客。"])
        assert ok is True
        user = s.get(User, alice_id)
        out = await runner._docs_context(
            [{"role": "user", "content": "主角叫什么名字?"}],
            [doc.id],
            user,
            s,
        )
    assert out is not None
    assert "世界观设定.md" in out
    assert "阿澈" in out


async def test_docs_context_ignores_other_users_doc(ctx):
    _, _, _, engine, (alice_id, bob_id) = ctx
    with Session(engine) as s:
        doc = Document(
            tenant_id="t", user_id=alice_id, filename="私密.txt", kind="txt", size=10
        )
        s.add(doc)
        s.commit()
        s.refresh(doc)
        await docsvc.build_index(alice_id, doc.id, ["秘密内容"])
        bob = s.get(User, bob_id)
        # bob 挂载 alice 的文档 id → 所有权过滤后无可检索文档 → None
        out = await runner._docs_context(
            [{"role": "user", "content": "秘密是什么?"}], [doc.id], bob, s
        )
    assert out is None


async def test_docs_context_empty_when_no_ids(ctx):
    _, _, _, engine, (alice_id, _) = ctx
    with Session(engine) as s:
        user = s.get(User, alice_id)
        out = await runner._docs_context(
            [{"role": "user", "content": "hi"}], [], user, s
        )
    assert out is None


async def test_retrieve_degrades_when_embedding_down(ctx, monkeypatch):
    """索引存在但查询时 embedding 挂了 → 检索返回空(不炸)。"""
    _, _, _, engine, (alice_id, _) = ctx
    with Session(engine) as s:
        doc = Document(
            tenant_id="t", user_id=alice_id, filename="a.txt", kind="txt", size=10
        )
        s.add(doc)
        s.commit()
        s.refresh(doc)
        await docsvc.build_index(alice_id, doc.id, ["一些内容"])

        async def down_embed(texts):
            return None

        monkeypatch.setattr(rag, "_embed", down_embed)
        hits = await docsvc.retrieve([doc], "内容")
        assert hits == []


# --------------------------------------------------------------------------- #
# 全格式扩展(2026-08-28):office/csv/json/代码/图片
# --------------------------------------------------------------------------- #
def test_kind_from_filename_full_formats():
    # office
    assert docsvc.kind_from_filename("报表.XLSX") == "xlsx"
    assert docsvc.kind_from_filename("路演.pptx") == "pptx"
    # 数据
    assert docsvc.kind_from_filename("数据.csv") == "csv"
    assert docsvc.kind_from_filename("配置.json") == "json"
    # 代码/文本
    assert docsvc.kind_from_filename("main.py") == "py"
    assert docsvc.kind_from_filename("app.tsx") == "tsx"
    assert docsvc.kind_from_filename("query.sql") == "sql"
    # 图片(文档通道走 VLM 反推)
    assert docsvc.kind_from_filename("照片.JPG") == "jpg"
    assert docsvc.kind_from_filename("a.png") == "png"
    assert docsvc.kind_from_filename("扫描.tiff") == "tiff"
    assert docsvc.kind_from_filename("a.bmp") == "bmp"
    # 仍然拒绝
    assert docsvc.kind_from_filename("a.exe") is None
    assert docsvc.kind_from_filename("a.zip") is None
    assert docsvc.is_image_kind("png") is True
    assert docsvc.is_image_kind("pdf") is False


def test_parse_text_csv_structured_preview():
    csv_body = "姓名,年龄,城市\n阿澈,27,杭州\n小满,31,上海\n"
    text = docsvc.parse_text("csv", csv_body.encode())
    assert "共 2 行数据、3 列" in text
    assert "列名:姓名, 年龄, 城市" in text
    assert "| 姓名 | 年龄 | 城市 |" in text
    assert "阿澈" in text


def test_parse_text_csv_long_table_truncates_sample():
    lines = ["id,value"] + [f"{i},v{i}" for i in range(docsvc._CSV_SAMPLE_ROWS + 10)]
    text = docsvc.parse_text("csv", "\n".join(lines).encode())
    assert f"共 {docsvc._CSV_SAMPLE_ROWS + 10} 行数据" in text
    assert "其余 10 行略" in text


def test_parse_text_json_object_summary_and_pretty():
    payload = {"name": "阿澈", "skills": ["剑", "酒"], "level": 9}
    text = docsvc.parse_text("json", json.dumps(payload).encode())
    assert "JSON 对象:共 3 个键" in text
    assert "name" in text and '"阿澈"' in text


def test_parse_text_json_array_summary():
    text = docsvc.parse_text("json", json.dumps([1, 2, 3]).encode())
    assert "JSON 数组:共 3 个元素" in text


def test_parse_text_json_invalid_raises():
    with pytest.raises(ValueError):
        docsvc.parse_text("json", b"{not json")


def test_parse_text_code_direct_read():
    src = "def hello():\n    return '世界'\n"
    assert docsvc.parse_text("py", src.encode()) == src


def test_parse_text_docx_with_table():
    import docx

    buf = io.BytesIO()
    d = docx.Document()
    d.add_paragraph("合同正文")
    table = d.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "条款"
    table.cell(0, 1).text = "金额"
    table.cell(1, 0).text = "首付"
    table.cell(1, 1).text = "5000"
    d.save(buf)
    text = docsvc.parse_text("docx", buf.getvalue())
    assert "合同正文" in text
    assert "| 条款 | 金额 |" in text and "首付" in text


def test_parse_text_xlsx_multi_sheet():
    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "收入"
    ws1.append(["月份", "金额"])
    ws1.append(["一月", 100])
    ws2 = wb.create_sheet("支出")
    ws2.append(["项目", "金额"])
    ws2.append(["房租", 50])
    buf = io.BytesIO()
    wb.save(buf)
    text = docsvc.parse_text("xlsx", buf.getvalue())
    assert "## Sheet:收入" in text and "| 月份 | 金额 |" in text and "一月" in text
    assert "## Sheet:支出" in text and "房租" in text


def test_parse_text_pptx_slides_and_notes():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "发布计划"
    slide.notes_slide.notes_text_frame.text = "内部口径"
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "预算"
    buf = io.BytesIO()
    prs.save(buf)
    text = docsvc.parse_text("pptx", buf.getvalue())
    assert "## 第 1 页" in text and "发布计划" in text
    assert "备注:内部口径" in text
    assert "## 第 2 页" in text and "预算" in text


def test_parse_text_pdf_scanned_marker():
    """无文本空白 PDF(扫描件语义)→ 返回标注占位而不是空串(路由不再 422)。"""
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    buf = io.BytesIO()
    writer.write(buf)
    text = docsvc.parse_text("pdf", buf.getvalue())
    assert text == docsvc.PDF_SCANNED_MARKER


def test_image_to_vlm_bytes_passthrough_and_convert():
    # jpg 原样直传
    raw = b"\xff\xd8\xff\xe0 fake"
    payload, mime = docsvc._image_to_vlm_bytes("jpg", raw)
    assert payload == raw and mime == "image/jpeg"
    # bmp → PNG 转码(PIL 真实往返)
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (4, 4), (255, 0, 0)).save(buf, format="BMP")
    payload, mime = docsvc._image_to_vlm_bytes("bmp", buf.getvalue())
    assert mime == "image/png" and payload.startswith(b"\x89PNG")


async def test_parse_document_routes_image_to_vlm(monkeypatch):
    called = {}

    async def fake_describe(kind, raw, filename):
        called["kind"] = kind
        return "【图片反推描述】一只橘猫"

    monkeypatch.setattr(docsvc, "describe_image", fake_describe)
    text = await docsvc.parse_document("png", b"\x89PNG fake", "猫.png")
    assert text == "【图片反推描述】一只橘猫"
    assert called["kind"] == "png"


async def test_parse_document_text_goes_local():
    text = await docsvc.parse_document("txt", "本地内容".encode(), "a.txt")
    assert text == "本地内容"


def test_upload_code_file_via_endpoint(ctx):
    c, alice_token, _, _, _ = ctx
    r = c.post(
        "/api/docs/upload",
        files={"file": ("脚本.py", "print('你好')\n".encode(), "text/x-python")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 201, r.text
    doc = r.json()
    assert doc["kind"] == "py" and doc["status"] == "ready"


def test_upload_xlsx_via_endpoint(ctx):
    import openpyxl

    c, alice_token, _, _, _ = ctx
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["指标", "数值"])
    ws.append(["销量", 42])
    buf = io.BytesIO()
    wb.save(buf)
    r = c.post(
        "/api/docs/upload",
        files={"file": ("报表.xlsx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 201, r.text
    assert r.json()["kind"] == "xlsx"
    # 切块内容含 markdown 表(经索引文件验证)
    doc_id = r.json()["id"]
    _, _, _, _, (alice_id, _) = ctx
    index = docsvc._index_path(alice_id, doc_id)
    data = json.loads(index.read_text(encoding="utf-8"))
    assert any("销量" in chunk for chunk in data["chunks"])


def test_upload_image_via_endpoint_uses_vlm(ctx, monkeypatch):
    c, alice_token, _, _, _ = ctx

    async def fake_describe(kind, raw, filename):
        return "【图片反推描述】海报上写着「夏日大促」"

    monkeypatch.setattr(docsvc, "describe_image", fake_describe)
    r = c.post(
        "/api/docs/upload",
        files={"file": ("海报.png", b"\x89PNG\r\n\x1a\n fake", "image/png")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 201, r.text
    doc = r.json()
    assert doc["kind"] == "png" and doc["status"] == "ready"


def test_upload_image_vlm_failure_passthrough_502(ctx, monkeypatch):
    from fastapi import HTTPException

    c, alice_token, _, _, _ = ctx

    async def down_describe(kind, raw, filename):
        raise HTTPException(502, "VLM 反推服务不可达")

    monkeypatch.setattr(docsvc, "describe_image", down_describe)
    r = c.post(
        "/api/docs/upload",
        files={"file": ("图.png", b"\x89PNG fake", "image/png")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 502  # 引擎不可达原码上抛,不伪装成 422 文件损坏


def test_upload_rejects_still_unsupported(ctx):
    c, alice_token, _, _, _ = ctx
    r = c.post(
        "/api/docs/upload",
        files={"file": ("压缩包.zip", b"PK\x03\x04", "application/zip")},
        headers=_auth(alice_token),
    )
    assert r.status_code == 400
