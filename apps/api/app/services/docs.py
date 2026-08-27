"""用户文档:全格式解析 → 切块 → embedding 索引 → 余弦检索。

格式矩阵(按扩展名路由,2026-08-28 全格式扩展):
- 直读文本/代码:txt/md/py/js/ts/html/css/sql/sh/yaml 等(_TEXT_KINDS)——utf-8 直读;
- pdf:pypdf 逐页提取;扫描件提取不到文本时返回标注占位(不拒绝上传,助手可如实告知);
- docx:python-docx 段落 + 表格(markdown);xlsx:openpyxl 逐 sheet 转 markdown 表;
  pptx:python-pptx 逐页提取文本框/表格/备注;
- csv/json:结构化预览(行列数/键结构摘要 + 头部样本),防大表格块爆炸又保留可检索性;
- 图片(jpg/png/webp/gif/bmp/tiff):复用 /api/reverse 的 VLM 反推链(spark01 Qwen3-VL)
  生成中文画面描述作为可检索文本,而非丢弃;tiff/bmp 先经 PIL 转 PNG 再送 VLM。

- 原文与索引(chunk+归一化向量 JSON)落盘 content_dir/docs/<user_id>/,不入库、不入 git。
- embedding 复用 agent.rag._embed(OpenAI 兼容端点,生产 Qwen3-Embedding-4B);
  不可用时文档仍保存,状态标注 no_embed,检索降级为空(与知识库 RAG 同一套优雅降级)。
- 长文档保护:chunk 上限 MAX_CHUNKS,超限截断并标注 status=partial,
  防止 500 页 PDF 的块数量爆炸拖垮 embedding。
"""
from __future__ import annotations

import asyncio
import csv
import io
import json
import logging
import re
from dataclasses import dataclass
from pathlib import Path

from app.agent import rag
from app.models import Document
from app.storage import content_subdir

logger = logging.getLogger(__name__)

MAX_FILE_BYTES = 50 * 1024 * 1024  # 单文件 ≤50MB
MAX_CHUNKS = 512  # 单文档 chunk 上限(防长文档块爆炸)
CHUNK_CHARS = 900  # 单块目标长度(字符)
CHUNK_OVERLAP = 120  # 超长段落硬切时的重叠

# 扫描件 PDF 占位标注:作为文档文本入索引,助手读到后能如实告知用户「不可提取文本」
PDF_SCANNED_MARKER = "(该 PDF 未能提取到文本内容,可能是扫描件或纯图片文档,暂不支持 OCR 识别)"

# 直读文本/代码扩展名(utf-8 errors=replace,永不炸)
_TEXT_KINDS = {
    "txt", "md", "markdown", "log", "csv", "json",
    "py", "js", "jsx", "ts", "tsx", "mjs", "cjs", "vue", "svelte",
    "html", "htm", "css", "scss", "less", "xml", "svg",
    "yaml", "yml", "toml", "ini", "cfg", "env",
    "sh", "bash", "zsh", "sql", "graphql",
    "java", "c", "h", "cpp", "hpp", "cc", "go", "rs", "rb", "php",
    "swift", "kt", "kts", "lua", "r", "scala", "pl", "ipynb",
}
_OFFICE_KINDS = {"pdf", "docx", "xlsx", "pptx"}
# 图片走 VLM 反推描述(文档通道传入的图片不再是「不支持的格式」)
_IMAGE_KINDS = {"jpg", "jpeg", "png", "webp", "gif", "bmp", "tiff", "tif"}
# VLM 不认的格式先经 PIL 转 PNG
_IMAGE_CONVERT_TO_PNG = {"bmp", "tiff", "tif"}

_KINDS = _TEXT_KINDS | _OFFICE_KINDS | _IMAGE_KINDS

# 结构化预览容量护栏
_CSV_SAMPLE_ROWS = 30  # csv 头部样本行数(超出标注「其余略」)
_XLSX_MAX_ROWS_PER_SHEET = 200  # 单 sheet 转 markdown 的行数上限
_JSON_PRETTY_CAP = 50_000  # json 美化输出字符上限
_TABLE_CELL_CAP = 80  # markdown 表单元格字符上限(防单元格内长文本撑爆块)


def kind_from_filename(filename: str) -> str | None:
    """按扩展名判定文档类型;不支持返回 None。"""
    suffix = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return suffix if suffix in _KINDS else None


def is_image_kind(kind: str) -> bool:
    return kind in _IMAGE_KINDS


def docs_root() -> Path:
    """文档存储根(content_dir/docs;不可写时 content_subdir 已降级临时目录)。"""
    return content_subdir("docs")


def _raw_path(user_id: str, doc_id: str, kind: str) -> Path:
    return docs_root() / user_id / f"{doc_id}.{kind}"


def _index_path(user_id: str, doc_id: str) -> Path:
    return docs_root() / user_id / f"{doc_id}.json"


# --------------------------------------------------------------------------- #
# 解析器:全部返回可检索纯文本/markdown
# --------------------------------------------------------------------------- #
def _md_cell(value: object) -> str:
    """单元格 → markdown 安全文本(去换行/转义管道/截长)。"""
    text = "" if value is None else str(value)
    text = re.sub(r"\s+", " ", text).replace("|", "\\|").strip()
    return text[: _TABLE_CELL_CAP - 1] + "…" if len(text) > _TABLE_CELL_CAP else text


def _md_table(rows: list[list[object]]) -> str:
    """行列表 → markdown 表(首行作表头;少于 1 行返回空)。"""
    if not rows:
        return ""
    header = [_md_cell(c) for c in rows[0]]
    lines = ["| " + " | ".join(header) + " |", "|" + "---|" * len(header)]
    for row in rows[1:]:
        cells = [_md_cell(c) for c in row]
        # 列数不齐时以表头为准补齐/截断,保证 markdown 表合法
        cells = (cells + [""] * len(header))[: len(header)]
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def _decode_text(raw: bytes) -> str:
    """utf-8(兼容 BOM)直读,坏字节 replace 容错。"""
    return raw.decode("utf-8-sig", errors="replace")


def _csv_text(raw: bytes) -> str:
    """csv → 结构化预览:行列数/列名摘要 + 头部样本 markdown 表。"""
    text = _decode_text(raw)
    rows = list(csv.reader(io.StringIO(text)))
    rows = [r for r in rows if any(cell.strip() for cell in r)]  # 去空行
    if not rows:
        return ""
    header, data = rows[0], rows[1:]
    summary = (
        f"CSV 表格:共 {len(data)} 行数据、{len(header)} 列。"
        f"列名:{', '.join(h.strip() for h in header)}"
    )
    sample = _md_table([header, *data[:_CSV_SAMPLE_ROWS]])
    suffix = f"\n(仅展示前 {_CSV_SAMPLE_ROWS} 行样本,其余 {len(data) - _CSV_SAMPLE_ROWS} 行略)" if len(data) > _CSV_SAMPLE_ROWS else ""
    return f"{summary}\n\n{sample}{suffix}"


def _json_text(raw: bytes) -> str:
    """json → 顶层结构摘要 + 美化样本(超大截断标注)。"""
    text = _decode_text(raw)
    obj = json.loads(text)  # 非法 JSON 上抛,由路由转 422
    if isinstance(obj, dict):
        keys = list(obj.keys())
        shown = ", ".join(str(k) for k in keys[:20])
        summary = f"JSON 对象:共 {len(keys)} 个键({shown}{', …' if len(keys) > 20 else ''})"
    elif isinstance(obj, list):
        summary = f"JSON 数组:共 {len(obj)} 个元素"
    else:
        summary = f"JSON 标量:{type(obj).__name__}"
    pretty = json.dumps(obj, ensure_ascii=False, indent=2)
    suffix = ""
    if len(pretty) > _JSON_PRETTY_CAP:
        pretty = pretty[:_JSON_PRETTY_CAP]
        suffix = f"\n(内容过长,仅展示前 {_JSON_PRETTY_CAP} 字符,其余略)"
    return f"{summary}\n\n{pretty}{suffix}"


def _pdf_text(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    # 扫描件/纯图片 PDF:不拒绝,落标注文本让助手如实告知(产品语义:可见即可问)
    return text if text.strip() else PDF_SCANNED_MARKER


def _docx_text(raw: bytes) -> str:
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = docx.Document(io.BytesIO(raw))
    parts: list[str] = []
    # iter_inner_content 保持段落/表格的文档顺序(python-docx ≥1.1)
    for block in document.iter_inner_content():
        if isinstance(block, Paragraph):
            if block.text.strip():
                parts.append(block.text)
        elif isinstance(block, Table):
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            table_md = _md_table(rows)
            if table_md:
                parts.append(table_md)
    return "\n".join(parts)


def _xlsx_text(raw: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        rows: list[list[object]] = []
        truncated = False
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= _XLSX_MAX_ROWS_PER_SHEET:
                truncated = True
                break
            rows.append(list(row))
        # 去掉全空尾行(read_only 模式 max_row 常虚高)
        while rows and not any(c is not None and str(c).strip() for c in rows[-1]):
            rows.pop()
        dims = f"{ws.max_row or 0} 行 × {ws.max_column or 0} 列"
        header = f"## Sheet:{ws.title}({dims})"
        table_md = _md_table(rows)
        note = f"\n(仅展示前 {_XLSX_MAX_ROWS_PER_SHEET} 行,其余略)" if truncated else ""
        parts.append(f"{header}\n{table_md}{note}" if table_md else header)
    wb.close()
    return "\n\n".join(parts)


def _pptx_text(raw: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
            elif shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table_md = _md_table(rows)
                if table_md:
                    texts.append(table_md)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            texts.append(f"(备注:{slide.notes_slide.notes_text_frame.text.strip()})")
        body = "\n".join(texts)
        if body:
            parts.append(f"## 第 {i} 页\n{body}")
    return "\n\n".join(parts)


def parse_text(kind: str, raw: bytes) -> str:
    """从原始字节解析纯文本(同步,CPU 密集;调用方负责 to_thread)。
    office 系惰性导入,缺依赖只影响对应格式。图片不走此函数(见 parse_document)。"""
    if kind in _TEXT_KINDS - {"csv", "json"}:
        return _decode_text(raw)
    if kind == "csv":
        return _csv_text(raw)
    if kind == "json":
        return _json_text(raw)
    if kind == "pdf":
        return _pdf_text(raw)
    if kind == "docx":
        return _docx_text(raw)
    if kind == "xlsx":
        return _xlsx_text(raw)
    if kind == "pptx":
        return _pptx_text(raw)
    raise ValueError(f"unsupported kind: {kind}")


def _image_to_vlm_bytes(kind: str, raw: bytes) -> tuple[bytes, str]:
    """VLM 不认的格式(tiff/bmp)经 PIL 转 PNG;其余原样直传。返回 (bytes, mime)。"""
    if kind not in _IMAGE_CONVERT_TO_PNG:
        mime = "image/jpeg" if kind in ("jpg", "jpeg") else f"image/{kind}"
        return raw, mime
    from PIL import Image

    with Image.open(io.BytesIO(raw)) as im:
        buf = io.BytesIO()
        im.convert("RGB").save(buf, format="PNG")
    return buf.getvalue(), "image/png"


async def describe_image(kind: str, raw: bytes, filename: str) -> str:
    """图片 → VLM 反推中文画面描述(复用 /api/reverse 的 Qwen3-VL 链路)。
    失败上抛 HTTPException,由路由转 502;不静默丢弃。"""
    from app.routes.reverse import reverse_visual  # 惰性导入防 routes↔services 环

    payload, mime = await asyncio.to_thread(_image_to_vlm_bytes, kind, raw)
    resp = await reverse_visual(payload, filename, mime, "image", False)
    return f"【图片反推描述】{resp.prompt}"


async def parse_document(kind: str, raw: bytes, filename: str) -> str:
    """统一异步解析入口:图片走 VLM 反推,其余格式本地解析(CPU 密集放线程池)。"""
    if is_image_kind(kind):
        return await describe_image(kind, raw, filename)
    return await asyncio.to_thread(parse_text, kind, raw)


# --------------------------------------------------------------------------- #
# 切块 / 落盘 / 索引 / 检索
# --------------------------------------------------------------------------- #
def split_chunks(text: str) -> tuple[list[str], bool]:
    """按段落聚合切块(≤CHUNK_CHARS),超长段落带重叠硬切。

    返回 (chunks, truncated):truncated=True 表示原文块数超 MAX_CHUNKS 被截断。
    """
    paras = [p.strip() for p in re.split(r"\n+", text) if p.strip()]
    chunks: list[str] = []
    buf = ""

    def flush() -> None:
        nonlocal buf
        if buf.strip():
            chunks.append(buf.strip())
        buf = ""

    for para in paras:
        while len(para) > CHUNK_CHARS:  # 超长段落(如无换行的 OCR 文本)硬切
            flush()
            chunks.append(para[:CHUNK_CHARS])
            para = para[CHUNK_CHARS - CHUNK_OVERLAP :]
        if buf and len(buf) + len(para) + 1 > CHUNK_CHARS:
            flush()
        buf = f"{buf}\n{para}" if buf else para
    flush()
    truncated = len(chunks) > MAX_CHUNKS
    return chunks[:MAX_CHUNKS], truncated


def save_raw(user_id: str, doc_id: str, kind: str, raw: bytes) -> None:
    """原文落盘(按用户分目录)。"""
    path = _raw_path(user_id, doc_id, kind)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(raw)


async def build_index(user_id: str, doc_id: str, chunks: list[str]) -> bool:
    """embedding 全部 chunk 并落盘 JSON 索引;embedding 不可用返回 False(不写索引)。"""
    vectors = await rag._embed(chunks)
    if vectors is None or len(vectors) != len(chunks):
        return False
    data = {
        "chunks": chunks,
        "vectors": [rag._normalize(v) for v in vectors],
    }
    path = _index_path(user_id, doc_id)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data), encoding="utf-8")
    return True


def delete_files(user_id: str, doc_id: str, kind: str) -> None:
    """删除原文与索引文件(不存在静默跳过)。"""
    for path in (_raw_path(user_id, doc_id, kind), _index_path(user_id, doc_id)):
        try:
            path.unlink(missing_ok=True)
        except OSError:
            logger.warning("删除文档文件失败: %s", path)


@dataclass(frozen=True)
class DocHit:
    filename: str
    text: str
    score: float


async def retrieve(
    docs: list[Document], query: str, k: int = 6, min_score: float = 0.25
) -> list[DocHit]:
    """对挂载文档做多文档余弦 top-k;无索引/embedding 不可用返回空(降级)。"""
    indexed = [d for d in docs if _index_path(d.user_id, d.id).is_file()]
    if not indexed or not query.strip():
        return []
    qe = await rag._embed([query])
    if not qe:
        return []
    q = rag._normalize(qe[0])
    hits: list[DocHit] = []
    for d in indexed:
        try:
            data = json.loads(_index_path(d.user_id, d.id).read_text(encoding="utf-8"))
        except (OSError, ValueError):
            continue
        chunks = data.get("chunks") or []
        vectors = data.get("vectors") or []
        for vec, text in zip(vectors, chunks):
            score = sum(a * b for a, b in zip(q, vec))
            if score >= min_score:
                hits.append(DocHit(filename=d.filename, text=text, score=score))
    hits.sort(key=lambda h: h.score, reverse=True)
    return hits[:k]
